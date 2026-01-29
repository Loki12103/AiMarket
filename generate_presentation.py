"""
AI Market Presentation Generator
Creates a professional PowerPoint presentation for the project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Powered Market Trend & Consumer Sentiment Forecaster"
    subtitle.text = "Lokesh K\nData Analyst & AI Developer\nSri Krishna College of Engineering and Technology\nJanuary 2026"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_introduction_slide(prs):
    """Slide 2: Introduction / Problem Statement"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Introduction & Problem Statement"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "The Challenge:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.level = 0
    
    points = [
        "Consumer opinions spread across social media, product reviews, and news platforms",
        "Businesses struggle to manually analyze massive, unstructured data in real time",
        "Delayed responses to market trends lead to customer dissatisfaction",
        "Need for automated, intelligent analysis of consumer sentiment"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.level = 1
    
    # Add solution
    p = tf.add_paragraph()
    p.text = "\nOur Solution:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AI-powered system that aggregates multi-source data and provides real-time sentiment and trend insights"
    p.font.size = Pt(18)
    p.level = 1

def create_objectives_slide(prs):
    """Slide 3: Objectives"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Project Objectives"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    objectives = [
        "Collect and aggregate consumer opinions from multiple channels",
        "Analyze sentiment and emerging topics using LLMs",
        "Provide contextual insights using RAG pipelines",
        "Help marketing and product teams make data-driven decisions",
        "Enable real-time monitoring and automated alerts"
    ]
    
    for i, obj in enumerate(objectives, 1):
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = obj
        p.font.size = Pt(20)
        p.level = 0
        p.space_after = Pt(12)

def create_architecture_slide(prs):
    """Slide 4: System Architecture"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Architecture layers
    layers = [
        ("Data Sources", "Social Media | Product Reviews | News APIs", 1.2),
        ("Collection Layer", "API Integration | Web Scraping | Data Aggregation", 2.2),
        ("Processing Layer", "Data Cleaning | Normalization | Text Processing", 3.2),
        ("AI Engine", "LLM Sentiment Analysis | Topic Modeling | Classification", 4.2),
        ("RAG Pipeline", "Vector Database (FAISS) | Semantic Search | Context Retrieval", 5.2),
        ("Output Layer", "Dashboards | Analytics | Alerts & Reports", 6.2)
    ]
    
    for layer_name, description, top in layers:
        # Layer box
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.5), Inches(top), Inches(7), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 102, 204)
        shape.line.color.rgb = RGBColor(0, 51, 102)
        
        # Layer text
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = layer_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        p = text_frame.add_paragraph()
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def create_modules_slide(prs):
    """Slide 5: Modules Overview"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Modules"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    modules = [
        ("1. Data Collection & Scraping", "Fetches data using APIs/scrapers from social media, reviews, and news"),
        ("2. Data Cleaning & Normalization", "Removes noise, formats text, standardizes dates and categories"),
        ("3. LLM Sentiment & Topic Modeling", "Classifies sentiment (Positive/Neutral/Negative) and extracts key topics"),
        ("4. RAG Pipeline", "Retrieves relevant data using vector search for contextual insights"),
        ("5. Visualization & Analytics", "Displays trends, comparisons, and sentiment shifts"),
        ("6. Alerts & Reporting", "Sends alerts for sentiment spikes and trend changes")
    ]
    
    for i, (module, desc) in enumerate(modules):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = module
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(8)

def create_milestones_slide(prs):
    """Slide 6: Project Milestones"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Project Milestones"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    milestones = [
        ("Milestone 1 (Week 2)", "Data scraping operational; cleaned dataset available"),
        ("Milestone 2 (Week 4)", "LLM sentiment & topic models functional"),
        ("Milestone 3 (Week 6)", "Dashboards live; RAG pipeline integrated"),
        ("Milestone 4 (Week 8)", "Alerts/reports generated; system deployed successfully")
    ]
    
    for i, (milestone, deliverable) in enumerate(milestones):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = milestone
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = deliverable
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(12)

def create_tech_stack_slide(prs):
    """Slide 7: Tech Stack"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech_categories = [
        ("Language", "Python 3.x"),
        ("Core Libraries", "pandas, numpy, scikit-learn"),
        ("AI/ML Libraries", "transformers, langchain, sentence-transformers"),
        ("LLM Models", "Gemini, Groq (Llama), BERT"),
        ("Topic Modeling", "LDA, BERTopic"),
        ("Vector Database", "FAISS"),
        ("Dashboard", "Streamlit, Plotly, Matplotlib"),
        ("Development Tools", "VS Code, Git, GitHub"),
        ("APIs", "RapidAPI, News API, Reddit API")
    ]
    
    for i, (category, tools) in enumerate(tech_categories):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{category}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 0
        
        run = p.runs[0]
        run2 = p.add_run()
        run2.text = tools
        run2.font.bold = False

def create_rag_llm_slide(prs):
    """Slide 8: RAG & LLM Integration"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "RAG & LLM Integration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "How it Works:"
    p.font.bold = True
    p.font.size = Pt(22)
    
    points = [
        "LLMs analyze sentiment and extract topics from consumer data",
        "RAG ensures insights are based on actual data, not hallucinations",
        "Vector search retrieves relevant opinions before generating insights",
        "FAISS enables fast semantic search across consumer feedback"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 1
    
    # Example query
    p = tf.add_paragraph()
    p.text = "\nExample Query:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.level = 0
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = '"Why is sentiment declining for Product X?"'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Response: "Negative sentiment increased by 25% due to delivery delays mentioned in 157 recent reviews."'
    p.font.size = Pt(16)
    p.level = 1

def create_results_slide(prs):
    """Slide 9: Results & Outputs"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Results & Key Outputs"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Deliverables:"
    p.font.bold = True
    p.font.size = Pt(22)
    
    results = [
        "✓ Interactive dashboards with real-time sentiment trends",
        "✓ Sentiment spike detection and visualization",
        "✓ Topic modeling with category-wise analysis",
        "✓ Automated alert system (Email/Slack notifications)",
        "✓ RAG-powered question answering system",
        "✓ Historical trend analysis and comparisons"
    ]
    
    for result in results:
        p = tf.add_paragraph()
        p.text = result
        p.font.size = Pt(18)
        p.level = 1
        p.space_after = Pt(8)
    
    # Example insight
    p = tf.add_paragraph()
    p.text = "\nExample Insight:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '"Negative sentiment increased 30% for Electronics category due to delivery delays and product quality issues"'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.level = 1

def create_testing_challenges_slide(prs):
    """Slide 10: Testing & Challenges"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Testing & Challenges"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Testing
    p = tf.paragraphs[0]
    p.text = "🧪 Testing:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 153, 76)
    
    testing = [
        "Tested with historical and real-time data",
        "Simulated sentiment spikes and trend changes",
        "Validated RAG pipeline accuracy",
        "Performance testing with large datasets"
    ]
    
    for test in testing:
        p = tf.add_paragraph()
        p.text = test
        p.font.size = Pt(16)
        p.level = 1
    
    # Challenges
    p = tf.add_paragraph()
    p.text = "\n⚠️ Challenges:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(204, 102, 0)
    p.level = 0
    
    challenges = [
        "Handling noisy social media text and slang",
        "Topic overlap in unsupervised models",
        "Sentiment imbalance across categories",
        "Performance optimization for large datasets",
        "API rate limits and data collection constraints"
    ]
    
    for challenge in challenges:
        p = tf.add_paragraph()
        p.text = challenge
        p.font.size = Pt(16)
        p.level = 1

def create_future_enhancements_slide(prs):
    """Slide 11: Future Enhancements"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Future Enhancements"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    enhancements = [
        "🌍 Multilingual sentiment analysis",
        "📈 Predictive trend forecasting using time series",
        "⚖️ Real-time brand comparison analytics",
        "🎤 Voice-of-customer analytics",
        "🔗 Integration with CRM tools (Salesforce, HubSpot)",
        "🤖 Advanced NLP models (GPT-4, Claude)",
        "📱 Mobile app for on-the-go insights",
        "🔔 Custom alert rules and thresholds"
    ]
    
    for i, enhancement in enumerate(enhancements):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = enhancement
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(10)

def create_conclusion_slide(prs):
    """Slide 12: Conclusion"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Key Takeaways:"
    p.font.bold = True
    p.font.size = Pt(24)
    
    points = [
        "Transforms unstructured consumer data into actionable intelligence",
        "Enables businesses to track sentiment and identify emerging trends",
        "Provides real-time, AI-driven insights for proactive decision-making",
        "Combines multiple data sources for comprehensive market understanding",
        "Reduces manual analysis time and improves response speed"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.level = 1
        p.space_after = Pt(10)
    
    # Impact
    p = tf.add_paragraph()
    p.text = "\nImpact:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.level = 0
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Empowers businesses to stay ahead of market changes and respond quickly to customer needs using cutting-edge AI technology."
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.level = 1

def create_thank_you_slide(prs):
    """Slide 13: Thank You"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Questions?\n\nLokesh K\nlokeshsh5252@gmail.com\nGitHub: github.com/Loki12103/AiMarket"
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.CENTER

def main():
    """Generate the complete presentation"""
    print("🎨 Generating AI Market Presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    create_title_slide(prs)
    print("✓ Title slide created")
    
    create_introduction_slide(prs)
    print("✓ Introduction slide created")
    
    create_objectives_slide(prs)
    print("✓ Objectives slide created")
    
    create_architecture_slide(prs)
    print("✓ Architecture slide created")
    
    create_modules_slide(prs)
    print("✓ Modules slide created")
    
    create_milestones_slide(prs)
    print("✓ Milestones slide created")
    
    create_tech_stack_slide(prs)
    print("✓ Tech stack slide created")
    
    create_rag_llm_slide(prs)
    print("✓ RAG & LLM slide created")
    
    create_results_slide(prs)
    print("✓ Results slide created")
    
    create_testing_challenges_slide(prs)
    print("✓ Testing & Challenges slide created")
    
    create_future_enhancements_slide(prs)
    print("✓ Future enhancements slide created")
    
    create_conclusion_slide(prs)
    print("✓ Conclusion slide created")
    
    create_thank_you_slide(prs)
    print("✓ Thank you slide created")
    
    # Save presentation
    filename = "AI_Market_Presentation.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation saved as: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n🎉 Ready to present!")

if __name__ == "__main__":
    main()
